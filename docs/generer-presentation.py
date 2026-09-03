# -*- coding: utf-8 -*-
"""Génère la présentation de soutenance de stage."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = r'c:\Users\ULTRAPC\Desktop\hcp'
GRAPHS = os.path.join(ROOT, 'graphs')
OUT = os.path.join(ROOT, 'docs', 'presentation-stage-HCP.pptx')

# ── palette sobre, dérivée de l'identité du HCP ──
VERT      = RGBColor(0x0B, 0x6E, 0x4F)
VERT_PALE = RGBColor(0xE8, 0xF2, 0xEE)
ENCRE     = RGBColor(0x14, 0x21, 0x1C)
GRIS      = RGBColor(0x53, 0x62, 0x5B)
GRIS_PALE = RGBColor(0xE2, 0xE6, 0xE4)
BLANC     = RGBColor(0xFF, 0xFF, 0xFF)
AMBRE     = RGBColor(0x8A, 0x5B, 0x08)

POLICE = 'Calibri'
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H


def _para(tf, texte, taille, gras=False, couleur=ENCRE, espace_avant=0,
          espace_apres=6, interligne=1.0, premier=False):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.text = texte
    p.space_before = Pt(espace_avant)
    p.space_after = Pt(espace_apres)
    p.line_spacing = interligne
    for r in p.runs:
        r.font.size = Pt(taille)
        r.font.bold = gras
        r.font.color.rgb = couleur
        r.font.name = POLICE
    return p


def zone(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def nouvelle(titre=None, sur_titre=None):
    """Une diapositive vierge, avec un titre et un filet vert."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if titre:
        if sur_titre:
            tf = zone(s, Inches(0.75), Inches(0.42), Inches(11.8), Inches(0.3))
            _para(tf, sur_titre.upper(), 11, True, VERT, espace_apres=0, premier=True)
            top = Inches(0.68)
        else:
            top = Inches(0.52)
        tf = zone(s, Inches(0.75), top, Inches(11.8), Inches(0.75))
        _para(tf, titre, 30, True, ENCRE, espace_apres=0, premier=True)
        # filet
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78),
                               top + Inches(0.62), Inches(1.5), Pt(3))
        r.fill.solid(); r.fill.fore_color.rgb = VERT
        r.line.fill.background(); r.shadow.inherit = False
    return s


def numeroter():
    for i, s in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        tf = zone(s, Inches(12.4), Inches(6.95), Inches(0.7), Inches(0.3))
        p = _para(tf, str(i), 10, False, GRIS, espace_apres=0, premier=True)
        p.alignment = PP_ALIGN.RIGHT


def image_ajustee(slide, nom, left, top, max_w, max_h):
    """Insère l'image en respectant ses proportions, centrée dans la boîte."""
    chemin = os.path.join(GRAPHS, nom)
    iw, ih = Image.open(chemin).size
    ratio = iw / ih
    if max_w / max_h > ratio:
        h = max_h; w = Emu(int(max_h * ratio))
    else:
        w = max_w; h = Emu(int(max_w / ratio))
    l = left + Emu(int((max_w - w) / 2))
    t = top + Emu(int((max_h - h) / 2))
    slide.shapes.add_picture(chemin, l, t, width=w, height=h)


def puces(slide, items, left, top, width, taille=16, hauteur=Inches(4.8)):
    tf = zone(slide, left, top, width, hauteur)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            titre, detail = it
            _para(tf, titre, taille, True, ENCRE, espace_avant=(0 if i == 0 else 10),
                  espace_apres=2, premier=(i == 0))
            _para(tf, detail, taille - 2, False, GRIS, espace_apres=0, interligne=1.15)
        else:
            _para(tf, it, taille, False, ENCRE, espace_avant=(0 if i == 0 else 6),
                  espace_apres=0, interligne=1.15, premier=(i == 0))
    return tf


def tableau(slide, entetes, lignes, left, top, width, largeurs=None, taille=13):
    nl, nc = len(lignes) + 1, len(entetes)
    h = Inches(0.42) * nl
    g = slide.shapes.add_table(nl, nc, left, top, width, h).table
    if largeurs:
        for i, lw in enumerate(largeurs):
            g.columns[i].width = lw
    for c, txt in enumerate(entetes):
        cell = g.cell(0, c)
        cell.text = txt
        cell.fill.solid(); cell.fill.fore_color.rgb = VERT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(taille); r.font.bold = True
                r.font.color.rgb = BLANC; r.font.name = POLICE
    for li, ligne in enumerate(lignes, start=1):
        for c, txt in enumerate(ligne):
            cell = g.cell(li, c)
            cell.text = str(txt)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLANC if li % 2 else RGBColor(0xF7, 0xFA, 0xF9)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(taille)
                    r.font.color.rgb = ENCRE if c == 0 else GRIS
                    r.font.bold = (c == 0)
                    r.font.name = POLICE
    return g


def encadre(slide, texte, left, top, width, height, fond=VERT_PALE, couleur=ENCRE, taille=15):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    b.fill.solid(); b.fill.fore_color.rgb = fond
    b.line.color.rgb = GRIS_PALE; b.line.width = Pt(0.75)
    b.shadow.inherit = False
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.18)
    _para(tf, texte, taille, False, couleur, espace_apres=0, interligne=1.2, premier=True)
    return b


def chiffres(slide, valeurs, top):
    """Une rangée de chiffres clés."""
    n = len(valeurs)
    marge, ecart = Inches(0.75), Inches(0.25)
    largeur = Emu(int((W - 2 * marge - ecart * (n - 1)) / n))
    for i, (val, lib) in enumerate(valeurs):
        l = marge + Emu(int(i * (largeur + ecart)))
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, top, largeur, Inches(1.25))
        b.fill.solid(); b.fill.fore_color.rgb = BLANC
        b.line.color.rgb = GRIS_PALE; b.line.width = Pt(1)
        b.shadow.inherit = False
        tf = b.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.12)
        p = _para(tf, val, 30, True, VERT, espace_apres=0, premier=True)
        p.alignment = PP_ALIGN.CENTER
        p = _para(tf, lib, 12, False, GRIS, espace_apres=0)
        p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════ 1 · COUVERTURE
s = prs.slides.add_slide(prs.slide_layouts[6])
bande = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
bande.fill.solid(); bande.fill.fore_color.rgb = VERT
bande.line.fill.background(); bande.shadow.inherit = False

tf = zone(s, Inches(1.1), Inches(1.6), Inches(11), Inches(0.4))
_para(tf, "HAUT-COMMISSARIAT AU PLAN  ·  DIRECTION DE LA STATISTIQUE", 13, True, VERT,
      espace_apres=0, premier=True)

tf = zone(s, Inches(1.1), Inches(2.15), Inches(11), Inches(1.8))
_para(tf, "Application de collecte d'enquêtes", 42, True, ENCRE, espace_apres=4, premier=True)
_para(tf, "pour les enquêteurs de terrain", 42, True, ENCRE, espace_apres=14)
_para(tf, "Architecture hors-ligne d'abord et génération assistée des questionnaires",
      18, False, GRIS)

r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.13), Inches(4.75), Inches(2.2), Pt(3))
r.fill.solid(); r.fill.fore_color.rgb = VERT
r.line.fill.background(); r.shadow.inherit = False

tf = zone(s, Inches(1.1), Inches(5.1), Inches(11), Inches(1.4))
_para(tf, "ABDELLAH ELMIR", 17, True, ENCRE, espace_apres=3, premier=True)
_para(tf, "Élève ingénieur — Data Engineering & Intelligence Artificielle", 14, False, GRIS,
      espace_apres=3)
_para(tf, "Soutenance de stage — septembre 2026", 14, False, GRIS)

# ══════════════════════════════ 2 · PROBLÈME
s = nouvelle("Le problème à résoudre", "Contexte")
puces(s, [
    ("Les enquêteurs travaillent là où le réseau n'existe pas",
     "Zones rurales, montagne, régions enclavées. Une application qui exige Internet\n"
     "pour fonctionner y est inutilisable."),
    ("La collecte papier coûte cher en temps et en fiabilité",
     "Double saisie, erreurs de recopie, plusieurs semaines entre le terrain et la donnée\n"
     "exploitable."),
    ("Chaque enquête a une forme différente",
     "Développer une application par enquête est impossible à maintenir."),
    ("Préparer un questionnaire numérique est long",
     "Transcrire à la main un document Word de 100 questions en format structuré\n"
     "prend des heures et introduit des erreurs."),
], Inches(0.75), Inches(1.75), Inches(11.8), taille=17)

encadre(s, "Objectif : une application utilisable une journée entière sans réseau, "
           "capable d'accueillir n'importe quelle enquête sans être reprogrammée.",
        Inches(0.75), Inches(5.9), Inches(11.8), Inches(0.85))

# ══════════════════════════════ 3 · UTILISATEURS
s = nouvelle("Deux profils, deux interfaces", "Utilisateurs")

b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.75),
                       Inches(5.7), Inches(4.1))
b.fill.solid(); b.fill.fore_color.rgb = BLANC
b.line.color.rgb = VERT; b.line.width = Pt(1.5); b.shadow.inherit = False
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25); tf.margin_right = Inches(0.3)
_para(tf, "L'ENQUÊTEUR", 12, True, VERT, espace_apres=6, premier=True)
_para(tf, "Application mobile", 22, True, ENCRE, espace_apres=12)
for t in ["Télécharge les questionnaires", "Interroge les ménages sur le terrain",
          "Travaille sans connexion pendant des jours",
          "Consulte l'aide de chaque question hors ligne",
          "Change la langue : français ou arabe",
          "Ses réponses remontent automatiquement"]:
    _para(tf, "—  " + t, 14, False, GRIS, espace_apres=5, interligne=1.1)

b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.75),
                       Inches(5.7), Inches(4.1))
b.fill.solid(); b.fill.fore_color.rgb = BLANC
b.line.color.rgb = VERT; b.line.width = Pt(1.5); b.shadow.inherit = False
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25); tf.margin_right = Inches(0.3)
_para(tf, "L'ADMINISTRATEUR", 12, True, VERT, espace_apres=6, premier=True)
_para(tf, "Outil web", 22, True, ENCRE, espace_apres=12)
for t in ["Importe un document Word ou texte", "Fait convertir le document en questionnaire",
          "Relit, corrige et publie l'enquête", "Gère les comptes des enquêteurs",
          "Suit la collecte sur un tableau de bord",
          "Exporte les réponses pour analyse"]:
    _para(tf, "—  " + t, 14, False, GRIS, espace_apres=5, interligne=1.1)

encadre(s, "Deux rôles seulement : enquêteur et administrateur. "
           "Une contrainte en base de données interdit toute autre valeur.",
        Inches(0.75), Inches(6.05), Inches(11.8), Inches(0.7))

# ══════════════════════════════ 4 · CAS D'UTILISATION
s = nouvelle("Diagramme de cas d'utilisation", "Modélisation")
image_ajustee(s, 'usecasehcp.png', Inches(0.6), Inches(1.5), Inches(7.6), Inches(5.5))
puces(s, [
    ("Acteurs principaux", "L'enquêteur et l'administrateur : les personnes\nqui déclenchent les actions."),
    ("Acteurs secondaires", "Le service d'IA et le serveur : des systèmes\nexternes que l'application appelle."),
    ("En jaune", "Tout ce qui fonctionne sans aucune connexion.\nC'est la majorité du travail de l'enquêteur."),
], Inches(8.5), Inches(1.9), Inches(4.2), taille=14)

# ══════════════════════════════ 5 · TECHNOLOGIES
s = nouvelle("Les technologies utilisées", "Choix techniques")
tableau(s,
    ["Composant", "Technologie", "Pourquoi ce choix"],
    [["Application mobile", "Flutter / Dart",
      "Un seul code source pour Android, iOS, Windows et le web"],
     ["Gestion d'état", "Riverpod",
      "Injection de dépendances et interface réactive sans code manuel"],
     ["Base locale", "SQLite via Drift",
      "Requêtes vérifiées à la compilation, résultats réactifs"],
     ["Communication", "API REST + JSON",
      "Standard, simple à déboguer, indépendant du client"],
     ["Serveur", "Node.js / Express",
      "Léger, adapté à une API de données"],
     ["Base centrale", "PostgreSQL",
      "Colonnes JSONB : accueille n'importe quelle forme d'enquête"],
     ["Conversion IA", "Google Gemini",
      "Sortie structurée imposée par un schéma JSON"],
     ["Lecture des .docx", "mammoth",
      "Extraction déterministe du texte et de la structure"]],
    Inches(0.75), Inches(1.75), Inches(11.8),
    largeurs=[Inches(2.5), Inches(2.5), Inches(6.8)], taille=13)

# ══════════════════════════════ 6 · ARCHITECTURE
s = nouvelle("L'architecture hors-ligne d'abord", "Principe fondateur")
puces(s, [
    ("La base locale est la source de vérité",
     "Aucun écran de l'application n'appelle le réseau. Les écrans lisent SQLite,\n"
     "et SQLite seul. Le réseau ne fait que remplir cette base ou la vider."),
    ("Conséquence directe",
     "L'application se comporte exactement de la même façon avec ou sans connexion.\n"
     "Il n'y a pas de code conditionnel « si en ligne » dans les écrans."),
    ("Le questionnaire n'est pas programmé, il est décrit",
     "Une enquête est un fichier JSON stocké tel quel. Publier une nouvelle enquête\n"
     "ne demande ni migration de base, ni recompilation de l'application."),
], Inches(0.75), Inches(1.75), Inches(11.8), taille=17)

encadre(s, "Un questionnaire de 100 questions d'une forme totalement nouvelle "
           "peut être mis en service sans modifier une seule ligne de code.",
        Inches(0.75), Inches(5.6), Inches(11.8), Inches(0.85))

# ══════════════════════════════ 7 · SYNCHRONISATION
s = nouvelle("Le parcours d'une réponse", "Synchronisation")
image_ajustee(s, 'synchcp.png', Inches(0.6), Inches(1.45), Inches(8.4), Inches(5.6))
puces(s, [
    ("Hors ligne", "La réponse est enregistrée à chaque saisie\ndans SQLite. Rien ne part sur le réseau."),
    ("À la validation", "Le statut passe de « brouillon » à\n« en attente d'envoi »."),
    ("Au retour du réseau", "Le moteur envoie automatiquement\ntout ce qui est en attente."),
    ("Si l'envoi échoue", "Nouvelle tentative, jusqu'à cinq essais."),
], Inches(9.2), Inches(1.9), Inches(3.6), taille=13)

# ══════════════════════════════ 8 · IDEMPOTENCE
s = nouvelle("Aucun doublon, même si le réseau coupe", "Point technique clé")
puces(s, [
    ("L'identifiant est généré sur le téléphone",
     "Chaque réponse reçoit un UUID au moment de sa création, hors ligne, avant même\n"
     "d'avoir vu le serveur. Cet identifiant est la clé primaire côté serveur."),
    ("Le serveur fait un INSERT ... ON CONFLICT DO UPDATE",
     "Réenvoyer la même réponse met à jour la ligne existante au lieu d'en créer une\n"
     "seconde. L'opération est idempotente."),
    ("Sans cela",
     "Si la connexion coupe pendant l'envoi, l'enquêteur réessaie et la réponse serait\n"
     "comptée deux fois dans les statistiques nationales."),
], Inches(0.75), Inches(1.75), Inches(11.8), taille=17)

encadre(s, "Vérifié : trois réponses envoyées, puis une réenvoyée à l'identique. "
           "La base en contient toujours trois.",
        Inches(0.75), Inches(5.7), Inches(11.8), Inches(0.8))

# ══════════════════════════════ 9 · QUESTIONNAIRES DYNAMIQUES
s = nouvelle("Un questionnaire est une donnée, pas du code", "Questionnaires dynamiques")

tf = zone(s, Inches(0.75), Inches(1.7), Inches(6.2), Inches(4.6))
_para(tf, "Description JSON d'une question", 14, True, VERT, espace_apres=8, premier=True)
for ligne in ['{', '  "id": "q_monthly_income",', '  "type": "number",',
              '  "label": { "fr": "Revenu mensuel",', '             "ar": "الدخل الشهري" },',
              '  "required": true,', '  "validation": { "min": 0 },',
              '  "visible_when": {', '     "question": "q_has_income",',
              '     "equals": "yes" }', '}']:
    p = _para(tf, ligne, 13, False, ENCRE, espace_apres=1, interligne=1.05)
    for r in p.runs:
        r.font.name = 'Consolas'

puces(s, [
    ("L'application lit cette description et construit l'écran",
     "Une seule fonction décide quel widget afficher selon le type."),
    ("Six types réellement affichés",
     "Texte, nombre, date, choix unique, liste déroulante, choix multiples."),
    ("Affichage conditionnel",
     "Une question peut n'apparaître que selon une réponse précédente,\n"
     "et n'est alors jamais exigée."),
    ("Bilingue par construction",
     "Chaque libellé porte sa version française et arabe."),
], Inches(7.3), Inches(1.8), Inches(5.3), taille=14)

# ══════════════════════════════ 10 · PIPELINE IA
s = nouvelle("Du document Word au questionnaire", "Génération assistée")
image_ajustee(s, 'generatesurvey.png', Inches(0.6), Inches(1.45), Inches(4.6), Inches(5.6))
puces(s, [
    ("1 · Extraction", "La bibliothèque mammoth lit le .docx. Aucune IA à ce stade.\n"
                       "La structure du document sert à écarter les titres de section."),
    ("2 · Découpage", "Le document est traité par lots de 25 questions."),
    ("3 · Conversion", "Le modèle reçoit une liste numérotée et doit rendre\n"
                       "exactement le même nombre de questions."),
    ("4 · Validation", "Le résultat est vérifié et réparé si nécessaire.\n"
                       "Chaque correction est affichée à l'administrateur."),
    ("5 · Relecture humaine", "Rien n'est enregistré avant que l'administrateur\n"
                              "ait relu et validé l'aperçu."),
], Inches(5.7), Inches(1.75), Inches(6.9), taille=14)

# ══════════════════════════════ 11 · FIABILITÉ IA
s = nouvelle("Comment garantir que l'IA ne se trompe pas", "Fiabilité")
puces(s, [
    ("Le modèle ne peut pas répondre en texte libre",
     "Un schéma JSON lui est imposé par l'API. La forme de la réponse est contrainte\n"
     "structurellement, pas seulement demandée."),
    ("Le prompt encadre strictement",
     "« Tu es un convertisseur strict, tu n'es pas un assistant. » Interdiction de fusionner,\n"
     "de supprimer ou d'inventer une question."),
    ("Le nombre de questions est vérifié",
     "Si un lot de 25 questions en rend 24, le lot est refait."),
    ("Le résultat est réparé et le rapport affiché",
     "Type inconnu ramené à du texte, options dupliquées supprimées. Chaque réparation\n"
     "est signalée, jamais silencieuse."),
], Inches(0.75), Inches(1.7), Inches(11.8), taille=16)

encadre(s, "Problème réel rencontré : sur un document de 100 questions répétitives, "
           "le modèle n'en avait rendu que 10 en les fusionnant, avec un JSON parfaitement "
           "valide. C'est ce qui a imposé la liste numérotée et la vérification du nombre. "
           "Après correction : 100 questions sur 100, et 150 sur 150.",
        Inches(0.75), Inches(5.75), Inches(11.8), Inches(1.1),
        fond=RGBColor(0xFB, 0xF0, 0xDC), couleur=AMBRE, taille=14)

# ══════════════════════════════ 12 · MESURES
s = nouvelle("Résultats mesurés", "Évaluation")
tf = zone(s, Inches(0.75), Inches(1.6), Inches(11.8), Inches(0.5))
_para(tf, "Un harnais d'évaluation compare la sortie générée à des réponses correctes "
          "écrites à la main, et attribue une note sur 10.", 15, False, GRIS,
      espace_apres=0, premier=True)

chiffres(s, [("9,8 / 10", "avec le modèle"), ("7,5 / 10", "règles seules"),
             ("5 s", "pour 10 questions"), ("73 s", "pour 150 questions"),
             ("52", "tests automatisés")], Inches(2.3))

tableau(s,
    ["Critère évalué", "Poids", "Ce qui est vérifié"],
    [["Extraction", "3 points", "Les vraies questions, sans les titres de section"],
     ["Type de question", "3 points", "Choix unique, multiple, nombre, date, texte"],
     ["Options", "2 points", "Les modalités sont présentes et pertinentes"],
     ["Obligatoire", "1 point", "« facultatif » est respecté"],
     ["Bilingue", "1 point", "Libellé et aide en français et en arabe"]],
    Inches(0.75), Inches(3.95), Inches(11.8),
    largeurs=[Inches(3.2), Inches(1.6), Inches(7.0)], taille=13)

# ══════════════════════════════ 13 · EXPLOITATION
s = nouvelle("Récupérer et suivre les données", "Exploitation")
puces(s, [
    ("Export CSV par enquête",
     "Une ligne par réponse, une colonne par question. Les codes internes sont\n"
     "retraduits en libellés lisibles."),
    ("Traçabilité complète",
     "Chaque ligne porte le matricule de l'enquêteur, sa région, et les deux dates :\n"
     "remplissage et synchronisation."),
    ("Ouverture directe dans Excel",
     "Encodage UTF-8 avec BOM et séparateur point-virgule, sans quoi Excel dégrade\n"
     "les accents et l'arabe. Une variante virgule existe pour R, Python ou SPSS."),
    ("Export également disponible en arabe",
     "Les en-têtes et les modalités sortent en arabe sur demande."),
    ("Tableau de bord",
     "Réponses collectées, enquêtes actives, activité par jour, par enquêteur\n"
     "et par région."),
], Inches(0.75), Inches(1.7), Inches(11.8), taille=16)

# ══════════════════════════════ 14 · QUALITÉ
s = nouvelle("Tests et qualité du code", "Vérification")
chiffres(s, [("31", "tests application"), ("21", "tests serveur"),
             ("1", "scénario bout en bout"), ("0", "erreur d'analyse statique")],
         Inches(1.75))
puces(s, [
    ("Ce que les tests couvrent",
     "La lecture des questionnaires, le moteur de synchronisation et ses trois scénarios,\n"
     "l'échappement du fichier d'export, et la validation de la sortie du modèle."),
    ("Trois tests de non-régression",
     "Écrits après avoir trouvé de vrais défauts, pour empêcher leur retour."),
    ("Un scénario complet rejouable",
     "Une commande rejoue tout le parcours : conversion, publication, travail hors ligne,\n"
     "synchronisation, absence de doublon, cloisonnement entre enquêteurs, export."),
], Inches(0.75), Inches(3.35), Inches(11.8), taille=16, hauteur=Inches(3.9))

# ══════════════════════════════ 15 · ÉTAT
s = nouvelle("État d'avancement", "Bilan")
tableau(s,
    ["Fonctionnalité", "État", "Précision"],
    [["Authentification en ligne et hors ligne", "Terminé", "Session renouvelée automatiquement"],
     ["Téléchargement des questionnaires", "Terminé", "Mise en cache locale"],
     ["Affichage dynamique", "Partiel", "Six types affichés sur vingt prévus"],
     ["Saisie et sauvegarde hors ligne", "Terminé", "Enregistrement à chaque réponse"],
     ["Synchronisation montante", "Terminé", "Réessais et absence de doublon"],
     ["Synchronisation descendante", "À faire", "Le serveur ne renvoie pas l'historique"],
     ["Gestion des conflits", "À faire", "Le dernier envoi écrase le précédent"],
     ["Interface bilingue français / arabe", "Terminé", "Y compris le sens de lecture"],
     ["Génération par IA", "Terminé", "Avec relecture avant publication"],
     ["Export et tableau de bord", "Terminé", "CSV compatible Excel"]],
    Inches(0.75), Inches(1.7), Inches(11.8),
    largeurs=[Inches(4.6), Inches(1.7), Inches(5.5)], taille=12.5)

# ══════════════════════════════ 16 · SUITE
s = nouvelle("Prochaines étapes", "Perspectives")
puces(s, [
    ("Synchronisation descendante",
     "Aujourd'hui les réponses ne montent que vers le serveur. Si un enquêteur change\n"
     "de téléphone, son historique local est perdu. C'est la priorité suivante."),
    ("Gestion des conflits",
     "Détecter qu'une même réponse a été modifiée des deux côtés, plutôt que de laisser\n"
     "le dernier envoi écraser le précédent."),
    ("Compilation Android",
     "Le code est déjà multiplateforme. La production de l'application installable\n"
     "ne demande qu'un environnement de compilation."),
    ("Types de questions supplémentaires",
     "Localisation GPS, photographie, signature, selon les besoins réels du HCP."),
    ("Durcissement pour la mise en production",
     "HTTPS, migrations de base versionnées, intégration continue."),
], Inches(0.75), Inches(1.7), Inches(11.8), taille=16)

# ══════════════════════════════ 17 · FIN
s = prs.slides.add_slide(prs.slide_layouts[6])
bande = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
bande.fill.solid(); bande.fill.fore_color.rgb = VERT
bande.line.fill.background(); bande.shadow.inherit = False

tf = zone(s, Inches(1.1), Inches(2.4), Inches(11), Inches(2.2))
_para(tf, "En résumé", 15, True, VERT, espace_apres=14, premier=True)
_para(tf, "Une application de terrain réellement utilisable sans réseau,", 26, True,
      ENCRE, espace_apres=4)
_para(tf, "capable d'accueillir n'importe quelle enquête sans reprogrammation,", 26, True,
      ENCRE, espace_apres=4)
_para(tf, "et un outil qui divise par cent le temps de préparation d'un questionnaire.",
      26, True, ENCRE, espace_apres=20)
_para(tf, "Merci de votre attention.", 17, False, GRIS)

numeroter()
prs.save(OUT)
print('Présentation générée :', OUT)
print('Diapositives :', len(prs.slides.__iter__.__self__._sldIdLst))
