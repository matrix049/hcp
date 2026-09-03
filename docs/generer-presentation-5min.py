# -*- coding: utf-8 -*-
"""Version courte : 7 diapositives pour une présentation de 5 minutes."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = r'c:\Users\ULTRAPC\Desktop\hcp'
GRAPHS = os.path.join(ROOT, 'graphs')
OUT = os.path.join(ROOT, 'docs', 'presentation-stage-HCP-5min.pptx')

VERT      = RGBColor(0x0B, 0x6E, 0x4F)
VERT_PALE = RGBColor(0xE8, 0xF2, 0xEE)
ENCRE     = RGBColor(0x14, 0x21, 0x1C)
GRIS      = RGBColor(0x53, 0x62, 0x5B)
GRIS_PALE = RGBColor(0xE2, 0xE6, 0xE4)
BLANC     = RGBColor(0xFF, 0xFF, 0xFF)
AMBRE     = RGBColor(0x8A, 0x5B, 0x08)
POLICE = 'Calibri'

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H


def _para(tf, texte, taille, gras=False, couleur=ENCRE, avant=0, apres=6,
          interligne=1.0, premier=False, police=POLICE):
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.text = texte
    p.space_before = Pt(avant); p.space_after = Pt(apres); p.line_spacing = interligne
    for r in p.runs:
        r.font.size = Pt(taille); r.font.bold = gras
        r.font.color.rgb = couleur; r.font.name = police
    return p


def zone(s, l, t, w, h):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    return tf


def nouvelle(titre, sur_titre=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if sur_titre:
        tf = zone(s, Inches(0.75), Inches(0.42), Inches(11.8), Inches(0.3))
        _para(tf, sur_titre.upper(), 11, True, VERT, apres=0, premier=True)
        top = Inches(0.68)
    else:
        top = Inches(0.52)
    tf = zone(s, Inches(0.75), top, Inches(11.8), Inches(0.75))
    _para(tf, titre, 30, True, ENCRE, apres=0, premier=True)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), top + Inches(0.62),
                           Inches(1.5), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = VERT
    r.line.fill.background(); r.shadow.inherit = False
    return s


def image_ajustee(s, nom, left, top, max_w, max_h):
    chemin = os.path.join(GRAPHS, nom)
    iw, ih = Image.open(chemin).size
    ratio = iw / ih
    if max_w / max_h > ratio:
        h = max_h; w = Emu(int(max_h * ratio))
    else:
        w = max_w; h = Emu(int(max_w / ratio))
    s.shapes.add_picture(chemin, left + Emu(int((max_w - w) / 2)),
                         top + Emu(int((max_h - h) / 2)), width=w, height=h)


def puces(s, items, left, top, width, taille=16, hauteur=Inches(4.6)):
    tf = zone(s, left, top, width, hauteur)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            t, d = it
            _para(tf, t, taille, True, ENCRE, avant=(0 if i == 0 else 10), apres=2,
                  premier=(i == 0))
            _para(tf, d, taille - 2, False, GRIS, apres=0, interligne=1.15)
        else:
            _para(tf, it, taille, False, ENCRE, avant=(0 if i == 0 else 6), apres=0,
                  interligne=1.15, premier=(i == 0))
    return tf


def encadre(s, texte, left, top, width, height, fond=VERT_PALE, couleur=ENCRE, taille=15):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    b.fill.solid(); b.fill.fore_color.rgb = fond
    b.line.color.rgb = GRIS_PALE; b.line.width = Pt(0.75); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.16); tf.margin_bottom = Inches(0.16)
    _para(tf, texte, taille, False, couleur, apres=0, interligne=1.2, premier=True)
    return b


def chiffres(s, valeurs, top, hauteur=Inches(1.2)):
    n = len(valeurs); marge, ecart = Inches(0.75), Inches(0.25)
    largeur = Emu(int((W - 2 * marge - ecart * (n - 1)) / n))
    for i, (val, lib) in enumerate(valeurs):
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               marge + Emu(int(i * (largeur + ecart))), top, largeur, hauteur)
        b.fill.solid(); b.fill.fore_color.rgb = BLANC
        b.line.color.rgb = GRIS_PALE; b.line.width = Pt(1); b.shadow.inherit = False
        tf = b.text_frame; tf.word_wrap = True; tf.margin_top = Inches(0.1)
        p = _para(tf, val, 28, True, VERT, apres=0, premier=True); p.alignment = PP_ALIGN.CENTER
        p = _para(tf, lib, 11.5, False, GRIS, apres=0); p.alignment = PP_ALIGN.CENTER


def minutage(s, texte):
    tf = zone(s, Inches(10.9), Inches(0.42), Inches(1.9), Inches(0.3))
    p = _para(tf, texte, 10.5, True, GRIS_PALE, apres=0, premier=True)
    p.alignment = PP_ALIGN.RIGHT


# ═══════════ 1 · COUVERTURE  (15 s)
s = prs.slides.add_slide(prs.slide_layouts[6])
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
b.fill.solid(); b.fill.fore_color.rgb = VERT
b.line.fill.background(); b.shadow.inherit = False

tf = zone(s, Inches(1.1), Inches(1.85), Inches(11), Inches(0.4))
_para(tf, "HAUT-COMMISSARIAT AU PLAN  ·  DIRECTION DE LA STATISTIQUE", 13, True, VERT,
      apres=0, premier=True)
tf = zone(s, Inches(1.1), Inches(2.4), Inches(11), Inches(1.8))
_para(tf, "Application de collecte d'enquêtes", 40, True, ENCRE, apres=4, premier=True)
_para(tf, "pour les enquêteurs de terrain", 40, True, ENCRE, apres=14)
_para(tf, "Fonctionne sans réseau · Questionnaires générés depuis un document Word",
      17, False, GRIS)
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.13), Inches(4.95), Inches(2.2), Pt(3))
r.fill.solid(); r.fill.fore_color.rgb = VERT
r.line.fill.background(); r.shadow.inherit = False
tf = zone(s, Inches(1.1), Inches(5.3), Inches(11), Inches(1.2))
_para(tf, "ABDELLAH ELMIR", 16, True, ENCRE, apres=3, premier=True)
_para(tf, "Élève ingénieur — Data Engineering & Intelligence Artificielle", 13.5, False,
      GRIS, apres=3)
_para(tf, "Soutenance de stage — septembre 2026", 13.5, False, GRIS)

# ═══════════ 2 · PROBLÈME ET RÉPONSE  (45 s)
s = nouvelle("Le problème, et la réponse apportée", "Contexte")
minutage(s, "45 s")

b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8),
                       Inches(5.7), Inches(3.5))
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xFB, 0xF0, 0xDC)
b.line.color.rgb = GRIS_PALE; b.line.width = Pt(0.75); b.shadow.inherit = False
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25); tf.margin_right = Inches(0.3)
_para(tf, "LE PROBLÈME", 12, True, AMBRE, apres=10, premier=True)
for t in ["Les enquêteurs travaillent là où il n'y a pas de réseau",
          "La collecte papier impose une double saisie et des semaines de délai",
          "Chaque enquête a une forme différente : impossible de coder une application par enquête",
          "Transcrire un questionnaire Word à la main prend des heures"]:
    _para(tf, "—  " + t, 14, False, ENCRE, apres=8, interligne=1.1)

b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.8),
                       Inches(5.7), Inches(3.5))
b.fill.solid(); b.fill.fore_color.rgb = VERT_PALE
b.line.color.rgb = VERT; b.line.width = Pt(1.25); b.shadow.inherit = False
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.25); tf.margin_right = Inches(0.3)
_para(tf, "LA RÉPONSE", 12, True, VERT, apres=10, premier=True)
for t in ["La base locale du téléphone est la source de vérité : l'application ignore le réseau",
          "Le questionnaire est décrit en JSON, pas programmé : nouvelle enquête sans recompiler",
          "Un document Word est converti automatiquement, puis relu par l'administrateur",
          "Les réponses remontent seules dès que la connexion revient"]:
    _para(tf, "—  " + t, 14, False, ENCRE, apres=8, interligne=1.1)

encadre(s, "Flutter et Riverpod pour l'application  ·  SQLite et Drift en local  ·  "
           "Node.js, Express et PostgreSQL côté serveur  ·  Google Gemini pour la conversion",
        Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.75), taille=14)

# ═══════════ 3 · CAS D'UTILISATION  (60 s)
s = nouvelle("Qui fait quoi", "Modélisation")
minutage(s, "60 s")
image_ajustee(s, 'usecasehcp.png', Inches(0.6), Inches(1.5), Inches(7.6), Inches(5.5))
puces(s, [
    ("Deux rôles seulement",
     "L'enquêteur sur le terrain,\nl'administrateur au bureau."),
    ("Deux systèmes externes",
     "Le service d'IA et le serveur :\nl'application les appelle,\nelle ne les contient pas."),
    ("En jaune",
     "Tout ce qui fonctionne sans\naucune connexion — c'est\nla quasi-totalité du travail\nde l'enquêteur."),
], Inches(8.5), Inches(2.0), Inches(4.2), taille=15)

# ═══════════ 4 · HORS LIGNE  (75 s)
s = nouvelle("Comment ça marche sans réseau", "Architecture")
minutage(s, "75 s")
image_ajustee(s, 'synchcp.png', Inches(0.5), Inches(1.5), Inches(8.3), Inches(4.5))
puces(s, [
    ("La donnée est écrite en local à chaque saisie",
     "Aucun écran n'appelle le réseau."),
    ("À la validation : « en attente d'envoi »"),
    ("Au retour du réseau, tout part automatiquement"),
], Inches(9.1), Inches(1.9), Inches(3.7), taille=14, hauteur=Inches(2.6))

encadre(s, "Point clé : l'identifiant de chaque réponse est généré sur le téléphone. "
           "Si le réseau coupe en plein envoi et que l'enquêteur réessaie, le serveur "
           "reconnaît le même identifiant et met à jour la ligne au lieu d'en créer une "
           "seconde. Aucun doublon dans les statistiques nationales.",
        Inches(0.75), Inches(6.15), Inches(11.8), Inches(1.0), taille=14)

# ═══════════ 5 · GÉNÉRATION IA  (75 s)
s = nouvelle("Du document Word au questionnaire", "Génération assistée")
minutage(s, "75 s")
image_ajustee(s, 'generatesurvey.png', Inches(0.6), Inches(1.5), Inches(4.3), Inches(5.4))
puces(s, [
    ("Le texte est extrait sans IA",
     "La structure du .docx sert à écarter les titres de section."),
    ("Le modèle reçoit une liste numérotée",
     "Il doit rendre exactement le même nombre de questions.\n"
     "Chaque lot de 25 est vérifié, et refait s'il en manque une."),
    ("Le résultat est validé puis réparé",
     "Type inconnu ramené à du texte, options dupliquées supprimées.\n"
     "Chaque correction est signalée, jamais silencieuse."),
    ("Rien n'est publié sans relecture humaine",
     "L'administrateur voit l'enquête comme l'enquêteur la verra,\n"
     "corrige à la main ou redemande à l'IA, puis publie."),
], Inches(5.4), Inches(1.8), Inches(7.2), taille=15)

encadre(s, "Sur un document de 100 questions répétitives, le modèle n'en avait rendu que 10 "
           "en les fusionnant — avec un JSON parfaitement valide. C'est ce qui a imposé la "
           "liste numérotée et la vérification du nombre. Après correction : 100 sur 100, 150 sur 150.",
        Inches(5.4), Inches(5.85), Inches(7.2), Inches(1.1),
        fond=RGBColor(0xFB, 0xF0, 0xDC), couleur=AMBRE, taille=13.5)

# ═══════════ 6 · RÉSULTATS  (60 s)
s = nouvelle("Résultats et état d'avancement", "Bilan")
minutage(s, "60 s")
chiffres(s, [("9,8 / 10", "qualité de conversion"), ("5 s", "pour 10 questions"),
             ("150", "questions sans perte"), ("52", "tests automatisés"),
             ("0", "doublon possible")], Inches(1.75))

b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.3),
                       Inches(5.7), Inches(3.0))
b.fill.solid(); b.fill.fore_color.rgb = VERT_PALE
b.line.color.rgb = VERT; b.line.width = Pt(1.25); b.shadow.inherit = False
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.22); tf.margin_right = Inches(0.3)
_para(tf, "TERMINÉ", 12, True, VERT, apres=9, premier=True)
for t in ["Authentification en ligne et hors ligne",
          "Téléchargement et remplissage hors ligne",
          "Synchronisation automatique, sans doublon",
          "Interface et questionnaires bilingues français / arabe",
          "Conversion Word vers questionnaire, avec relecture",
          "Tableau de bord et export CSV pour les statisticiens"]:
    _para(tf, "—  " + t, 13.5, False, ENCRE, apres=6, interligne=1.1)

b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(3.3),
                       Inches(5.7), Inches(3.0))
b.fill.solid(); b.fill.fore_color.rgb = BLANC
b.line.color.rgb = GRIS_PALE; b.line.width = Pt(1.25); b.shadow.inherit = False
tf = b.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.22); tf.margin_right = Inches(0.3)
_para(tf, "RESTE À FAIRE", 12, True, AMBRE, apres=9, premier=True)
for t in ["Synchronisation descendante : le serveur ne renvoie pas encore l'historique",
          "Gestion des conflits : le dernier envoi écrase le précédent",
          "Quatorze types de questions restent à implémenter",
          "Compilation de l'application Android installable",
          "Durcissement pour la mise en production"]:
    _para(tf, "—  " + t, 13.5, False, GRIS, apres=6, interligne=1.1)

# ═══════════ 7 · CONCLUSION  (25 s)
s = prs.slides.add_slide(prs.slide_layouts[6])
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
b.fill.solid(); b.fill.fore_color.rgb = VERT
b.line.fill.background(); b.shadow.inherit = False
tf = zone(s, Inches(1.1), Inches(2.5), Inches(11), Inches(2.4))
_para(tf, "En résumé", 15, True, VERT, apres=16, premier=True)
_para(tf, "Une application de terrain réellement utilisable sans réseau,", 25, True,
      ENCRE, apres=5)
_para(tf, "capable d'accueillir n'importe quelle enquête sans reprogrammation,", 25, True,
      ENCRE, apres=5)
_para(tf, "et un questionnaire préparé en quelques secondes au lieu de plusieurs heures.",
      25, True, ENCRE, apres=22)
_para(tf, "Merci de votre attention.", 16, False, GRIS)

# numérotation
for i, sl in enumerate(prs.slides, start=1):
    if i == 1:
        continue
    tf = zone(sl, Inches(12.4), Inches(6.95), Inches(0.7), Inches(0.3))
    p = _para(tf, str(i), 10, False, GRIS, apres=0, premier=True)
    p.alignment = PP_ALIGN.RIGHT

prs.save(OUT)
print('Version courte generee :', OUT)
